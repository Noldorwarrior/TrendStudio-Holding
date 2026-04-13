"""
test_26_presentation_build.py
Тесты для B2 — презентации ТрендСтудио v1.4.3.

Проверяет:
* корректность pptx (12 слайдов, 16:9, якорь, метрики)
* корректность интерактивного HTML (навигация, Plotly, 12 слайдов)
* согласованность значений со stress_matrix артефактами
"""
from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

PIPELINE = Path(__file__).resolve().parent.parent
STRESS = PIPELINE / "artifacts" / "stress_matrix"
ARTIFACTS = PIPELINE / "artifacts"

PPTX_PATH = ARTIFACTS / "B2_presentation.pptx"
HTML_PATH = ARTIFACTS / "B2_presentation.html"


@pytest.fixture(scope="module")
def pptx_doc():
    assert PPTX_PATH.exists(), f"Missing {PPTX_PATH}"
    return Presentation(str(PPTX_PATH))


@pytest.fixture(scope="module")
def html_text():
    assert HTML_PATH.exists(), f"Missing {HTML_PATH}"
    return HTML_PATH.read_text(encoding="utf-8")


@pytest.fixture(scope="module")
def lhs_data():
    return json.loads((STRESS / "lhs_copula.json").read_text(encoding="utf-8"))


@pytest.fixture(scope="module")
def mc_data():
    return json.loads((STRESS / "monte_carlo.json").read_text(encoding="utf-8"))


# ---------- PPTX structure ----------
def test_pptx_exists(pptx_doc):
    assert pptx_doc is not None


def test_pptx_has_12_slides(pptx_doc):
    assert len(pptx_doc.slides) == 12, f"Ожидалось 12 слайдов, получено {len(pptx_doc.slides)}"


def test_pptx_is_16_9(pptx_doc):
    # 16:9 ratio — ширина к высоте близка к 1.777
    ratio = pptx_doc.slide_width / pptx_doc.slide_height
    assert 1.7 < ratio < 1.85, f"Ожидался 16:9, ratio={ratio}"


def test_pptx_slides_contain_text(pptx_doc):
    """Каждый слайд должен содержать хотя бы один текстовый фрейм."""
    for i, slide in enumerate(pptx_doc.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            texts.append(run.text)
        assert texts, f"Слайд {i} пустой"


def _all_text(pptx_doc) -> str:
    buf = []
    for slide in pptx_doc.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                buf.append(shape.text_frame.text)
    return "\n".join(buf)


def test_pptx_mentions_trend_studio(pptx_doc):
    assert "ТрендСтудио" in _all_text(pptx_doc)


def test_pptx_mentions_anchor_3000(pptx_doc):
    assert "3 000" in _all_text(pptx_doc) or "3000" in _all_text(pptx_doc)


def test_pptx_mentions_lhs_copula(pptx_doc):
    txt = _all_text(pptx_doc)
    assert "LHS" in txt and "Copula" in txt


def test_pptx_mentions_all_four_engines(pptx_doc):
    txt = _all_text(pptx_doc)
    assert "Naive MC" in txt
    assert "Bootstrap" in txt
    assert "Stage-Gate" in txt or "Stage-gate" in txt
    assert "LHS" in txt


def test_pptx_has_version(pptx_doc):
    assert "v1.4.3" in _all_text(pptx_doc)


def test_pptx_mentions_var95(pptx_doc):
    assert "VaR95" in _all_text(pptx_doc)


def test_pptx_mentions_tornado(pptx_doc):
    txt = _all_text(pptx_doc).lower()
    assert "tornado" in txt or "торнадо" in txt


def test_pptx_mentions_stage_gate_probability(pptx_doc, lhs_data):
    """72,1% — P(reach release)."""
    txt = _all_text(pptx_doc)
    assert "72,1" in txt or "72.1" in txt or "72%" in txt


# ---------- HTML structure ----------
def test_html_exists(html_text):
    assert html_text is not None and len(html_text) > 1000


def test_html_has_plotly(html_text):
    assert "plotly" in html_text.lower()


def test_html_has_12_slides(html_text):
    # Количество слайд-блоков
    assert html_text.count('class="slide') >= 12


def test_html_has_keyboard_navigation(html_text):
    assert "ArrowRight" in html_text and "ArrowLeft" in html_text


def test_html_has_swipe_navigation(html_text):
    assert "touchstart" in html_text and "touchend" in html_text


def test_html_has_progress_bar(html_text):
    assert 'id="progress"' in html_text
    assert 'class="progress"' in html_text


def test_html_has_slide_counter(html_text):
    assert 'id="counter"' in html_text


def test_html_has_chart_containers(html_text):
    assert 'id="chart-anchor"' in html_text
    assert 'id="chart-mc"' in html_text
    assert 'id="chart-tornado"' in html_text
    assert 'id="chart-funnel"' in html_text


def test_html_has_gauge_plot(html_text):
    assert "'indicator'" in html_text or '"indicator"' in html_text
    assert "gauge+number+delta" in html_text


def test_html_has_funnel_plot(html_text):
    assert "'funnel'" in html_text or '"funnel"' in html_text


def test_html_has_anchor_3000(html_text):
    assert "3000" in html_text


def test_html_has_animation_transitions(html_text):
    assert "transition" in html_text.lower()


def test_html_has_version(html_text):
    assert "v1.4.3" in html_text


# ---------- Cross-consistency ----------
def test_html_lhs_mean_matches_artifact(html_text, lhs_data):
    """LHS mean из HTML должен совпадать с lhs_copula.json."""
    assert f"{lhs_data['mean_ebitda']}" in html_text


def test_html_lhs_var95_matches(html_text, lhs_data):
    expected = f"{int(round(lhs_data['var_95_mln_rub']))}"
    assert expected in html_text or f"{lhs_data['var_95_mln_rub']}" in html_text


def test_html_has_four_engines(html_text):
    assert "Naive MC" in html_text
    assert "Bootstrap" in html_text or "Block" in html_text
    assert "Stage-Gate" in html_text or "Stage-gate" in html_text
    assert "LHS" in html_text and "Copula" in html_text


def test_html_has_roadmap(html_text):
    assert "Road-map" in html_text or "roadmap" in html_text.lower()


def test_html_has_risks(html_text):
    assert "риск" in html_text.lower()


def test_html_has_mitigation(html_text):
    assert "FX-хеджирование" in html_text or "хеджирование" in html_text.lower()


def test_html_uses_tbank_palette(html_text):
    # yellow accent #FFDD2D характерен для Т-банк палитры
    assert "#FFDD2D" in html_text or "FFDD2D" in html_text
